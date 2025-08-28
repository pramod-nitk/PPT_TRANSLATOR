# -*- coding: utf-8 -*-
"""
PPT Translator Module

A module for translating PowerPoint presentations using Google Translate.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from langcodes import Language
from googletrans import Translator
import time
import os


class PPTTranslator:
    """
    A class for translating PowerPoint presentations to multiple languages.
    """
    
    def __init__(self):
        """Initialize the translator with Google Translate service."""
        self.translator = Translator()
        self.rtl_langs = {"ar", "fa", "ur", "he"}
    
    def normalize_lang(self, lang: str) -> str:
        """
        Normalize language code to standard format.
        
        Args:
            lang (str): Language code or name
            
        Returns:
            str: Normalized language code
        """
        try:
            return Language.get(lang).to_tag().lower()
        except Exception:
            return lang.lower()
    
    def set_textframe_autofit(self, tf):
        """
        Set text frame to auto-fit content.
        
        Args:
            tf: Text frame object
        """
        try:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
    
    def iter_all_text_objects(self, shapes):
        """
        Iterate through all text objects in shapes.
        
        Args:
            shapes: Collection of shapes
            
        Yields:
            tuple: (kind, object) where kind is "text_frame" or "table_cell"
        """
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self.iter_all_text_objects(shape.shapes)
                continue

            if getattr(shape, "has_text_frame", False):
                yield ("text_frame", shape.text_frame)

            if getattr(shape, "has_table", False):
                table = shape.table
                for r in table.rows:
                    for c in r.cells:
                        if c.text_frame:
                            yield ("table_cell", c)
    
    def collect_runs_in_textframe(self, tf):
        """
        Collect all text runs in a text frame.
        
        Args:
            tf: Text frame object
            
        Returns:
            list: List of (paragraph, run) tuples
        """
        pairs = []
        for p in tf.paragraphs:
            for r in p.runs:
                pairs.append((p, r))
        return pairs
    
    def safe_translate(self, text, dest="fr", retries=3, delay=2):
        """
        Safely translate text with retry mechanism.
        
        Args:
            text (str): Text to translate
            dest (str): Destination language code
            retries (int): Number of retry attempts
            delay (int): Delay between retries in seconds
            
        Returns:
            str: Translated text or original text if translation fails
        """
        if not text.strip():
            return text
        
        for i in range(retries):
            try:
                result = self.translator.translate(text, dest=dest)
                if result and result.text:
                    return result.text
            except Exception as e:
                print(f"Translation error: {e} (retry {i+1}/{retries})")
                time.sleep(delay)
        
        return text
    
    def translate_batch(self, texts, target_lang):
        """
        Translate a batch of texts.
        
        Args:
            texts (list): List of texts to translate
            target_lang (str): Target language code
            
        Returns:
            list: List of translated texts
        """
        return [self.safe_translate(t, dest=target_lang) for t in texts]
    
    def translate_textframe(self, tf, target_lang_code: str):
        """
        Translate all text in a text frame.
        
        Args:
            tf: Text frame object
            target_lang_code (str): Target language code
        """
        self.set_textframe_autofit(tf)
        pairs = self.collect_runs_in_textframe(tf)
        original_texts = [r.text or "" for (_, r) in pairs]

        translated = self.translate_batch(original_texts, target_lang_code)

        for (p, r), new_text in zip(pairs, translated):
            r.text = new_text
    
    def translate_presentation(self, input_path: str, output_path: str, target_lang_code: str):
        """
        Translate an entire PowerPoint presentation.
        
        Args:
            input_path (str): Path to input PowerPoint file
            output_path (str): Path to save translated PowerPoint file
            target_lang_code (str): Target language code
            
        Raises:
            Exception: If translation fails
        """
        try:
            # Load presentation
            prs = Presentation(input_path)
            
            # Process each slide
            for slide in prs.slides:
                for kind, obj in self.iter_all_text_objects(slide.shapes):
                    if kind == "text_frame":
                        self.translate_textframe(obj, target_lang_code)
                    elif kind == "table_cell":
                        self.translate_textframe(obj.text_frame, target_lang_code)
            
            # Save translated presentation
            prs.save(output_path)
            
        except Exception as e:
            raise Exception(f"Failed to translate presentation: {str(e)}")
    
    def translate_multiple_languages(self, input_path: str, output_dir: str, languages: list):
        """
        Translate a presentation to multiple languages.
        
        Args:
            input_path (str): Path to input PowerPoint file
            output_dir (str): Directory to save translated files
            languages (list): List of language codes to translate to
            
        Returns:
            list: List of successfully translated file paths
        """
        os.makedirs(output_dir, exist_ok=True)
        translated_files = []
        
        for lang in languages:
            try:
                safe_lang = self.normalize_lang(lang).replace("/", "-")
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                out_path = os.path.join(output_dir, f"{base_name}_{safe_lang}.pptx")
                
                self.translate_presentation(input_path, out_path, safe_lang)
                translated_files.append(out_path)
                print(f"✔ Saved: {out_path}")
                
            except Exception as e:
                print(f"❌ Failed to translate to {lang}: {str(e)}")
        
        return translated_files
    
    def get_supported_languages(self):
        """
        Get list of supported languages.
        
        Returns:
            dict: Dictionary mapping language names to codes
        """
        return {
            "French": "fr",
            "German": "de", 
            "Spanish": "es",
            "Hindi": "hi",
            "Japanese": "ja",
            "Chinese (Simplified)": "zh-cn",
            "Russian": "ru",
            "Arabic": "ar",
            "Portuguese": "pt",
            "Italian": "it",
            "Korean": "ko",
            "Dutch": "nl",
            "Swedish": "sv",
            "Norwegian": "no",
            "Danish": "da",
            "Finnish": "fi",
            "Polish": "pl",
            "Turkish": "tr",
            "Greek": "el",
            "Hebrew": "he"
        }
